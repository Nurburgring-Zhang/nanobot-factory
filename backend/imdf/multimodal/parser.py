"""P4-7-W1: MultiModalParser — 文档/PDF/Office/Email/图像/视频/语音 统一解析.

Supported input formats
-----------------------

* 文档: PDF / Word / Excel / PowerPoint / Markdown / HTML / 纯文本
* 图像: JPEG / PNG / WEBP / GIF / BMP
* 视频: MP4 / AVI / MOV / MKV (关键帧提取)
* 音频: MP3 / WAV / FLAC (ASR via whisper if available)
* 邮件: EML / MSG / MBOX (header + body + 附件)

Every parsed artefact is normalised to a :class:`MultimodalDocument` with
the following shape (Google Flow Agent / Gemini-Omni style)::

    MultimodalDocument(
        modality="document",          # one of 6 modalities
        text="...extracted text...",
        segments=[DocumentSegment(...)],  # ordered text chunks
        images=[DocumentImage(...)],      # embedded images (base64 + meta)
        tables=[DocumentTable(...)],      # structured tables
        metadata={"source": ..., "pages": N, "duration": T, ...},
    )

The parser is best-effort: it never raises for missing optional
dependencies, instead returning a degraded document with a note in
metadata["warnings"].
"""
from __future__ import annotations

import base64
import email
import email.policy
import hashlib
import io
import json
import logging
import mimetypes
import os
import re
import struct
import tempfile
import uuid
import wave
import zipfile
from dataclasses import dataclass, field, asdict
from typing import Any, Dict, Iterable, List, Optional, Tuple, Union

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Public modality constants — shared with adapter + 12 services.
# ---------------------------------------------------------------------------
MODALITY_TEXT = "text"
MODALITY_IMAGE = "image"
MODALITY_VIDEO = "video"
MODALITY_AUDIO = "audio"
MODALITY_DOCUMENT = "document"
MODALITY_EMAIL = "email"
MODALITY_MULTIMIX = "multimodal_mix"

ALL_MODALITIES = (
    MODALITY_TEXT,
    MODALITY_IMAGE,
    MODALITY_VIDEO,
    MODALITY_AUDIO,
    MODALITY_DOCUMENT,
    MODALITY_EMAIL,
    MODALITY_MULTIMIX,
)

# Output kinds the 12 services expose
OUTPUT_TEXT = "text"
OUTPUT_JSON = "structured_json"
OUTPUT_MULTIMODAL = "multimodal_response"

ALL_OUTPUT_KINDS = (OUTPUT_TEXT, OUTPUT_JSON, OUTPUT_MULTIMODAL)

# Mapping from MIME / extension → canonical modality
MIME_TO_MODALITY: Dict[str, str] = {
    # text
    "text/plain": MODALITY_TEXT,
    "text/markdown": MODALITY_DOCUMENT,
    "text/html": MODALITY_DOCUMENT,
    # image
    "image/jpeg": MODALITY_IMAGE,
    "image/png": MODALITY_IMAGE,
    "image/webp": MODALITY_IMAGE,
    "image/gif": MODALITY_IMAGE,
    "image/bmp": MODALITY_IMAGE,
    # video
    "video/mp4": MODALITY_VIDEO,
    "video/avi": MODALITY_VIDEO,
    "video/quicktime": MODALITY_VIDEO,
    "video/x-matroska": MODALITY_VIDEO,
    "video/webm": MODALITY_VIDEO,
    # audio
    "audio/mpeg": MODALITY_AUDIO,
    "audio/mp3": MODALITY_AUDIO,
    "audio/wav": MODALITY_AUDIO,
    "audio/x-wav": MODALITY_AUDIO,
    "audio/flac": MODALITY_AUDIO,
    "audio/x-m4a": MODALITY_AUDIO,
    # document
    "application/pdf": MODALITY_DOCUMENT,
    "application/msword": MODALITY_DOCUMENT,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": MODALITY_DOCUMENT,
    "application/vnd.ms-excel": MODALITY_DOCUMENT,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": MODALITY_DOCUMENT,
    "application/vnd.ms-powerpoint": MODALITY_DOCUMENT,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": MODALITY_DOCUMENT,
    # email
    "message/rfc822": MODALITY_EMAIL,
    "application/vnd.ms-outlook": MODALITY_EMAIL,
    "application/mbox": MODALITY_EMAIL,
}

EXT_TO_MODALITY: Dict[str, str] = {
    ".txt": MODALITY_TEXT,
    ".md": MODALITY_DOCUMENT, ".markdown": MODALITY_DOCUMENT,
    ".html": MODALITY_DOCUMENT, ".htm": MODALITY_DOCUMENT,
    ".pdf": MODALITY_DOCUMENT,
    ".doc": MODALITY_DOCUMENT, ".docx": MODALITY_DOCUMENT,
    ".xls": MODALITY_DOCUMENT, ".xlsx": MODALITY_DOCUMENT,
    ".ppt": MODALITY_DOCUMENT, ".pptx": MODALITY_DOCUMENT,
    ".jpg": MODALITY_IMAGE, ".jpeg": MODALITY_IMAGE,
    ".png": MODALITY_IMAGE, ".webp": MODALITY_IMAGE,
    ".gif": MODALITY_IMAGE, ".bmp": MODALITY_IMAGE,
    ".mp4": MODALITY_VIDEO, ".avi": MODALITY_VIDEO,
    ".mov": MODALITY_VIDEO, ".mkv": MODALITY_VIDEO, ".webm": MODALITY_VIDEO,
    ".mp3": MODALITY_AUDIO, ".wav": MODALITY_AUDIO, ".flac": MODALITY_AUDIO,
    ".m4a": MODALITY_AUDIO, ".ogg": MODALITY_AUDIO,
    ".eml": MODALITY_EMAIL, ".msg": MODALITY_EMAIL, ".mbox": MODALITY_EMAIL,
}


# ---------------------------------------------------------------------------
# Data classes
# ---------------------------------------------------------------------------
@dataclass
class DocumentImage:
    """An image embedded inside a parsed document (or a raw image input)."""
    image_id: str
    mime_type: str
    width: int = 0
    height: int = 0
    page: Optional[int] = None
    timestamp: Optional[float] = None  # for video frames
    ocr_text: str = ""
    base64: str = ""  # populated only when explicitly requested
    bytes_size: int = 0
    sha256: str = ""
    caption: str = ""  # alt text / OCR caption

    def to_dict(self) -> Dict[str, Any]:
        d = asdict(self)
        # drop base64 unless small; most consumers want the meta not the blob
        if len(d.get("base64") or "") > 256:
            d["base64"] = f"<{len(d['base64'])//4} bytes>"
        return d


@dataclass
class DocumentTable:
    """Structured table (rows of cells)."""
    table_id: str
    rows: List[List[str]] = field(default_factory=list)
    page: Optional[int] = None
    caption: str = ""

    def to_dict(self) -> Dict[str, Any]:
        return asdict(self)

    @property
    def n_rows(self) -> int:
        return len(self.rows)

    @property
    def n_cols(self) -> int:
        if not self.rows:
            return 0
        return max(len(r) for r in self.rows)


@dataclass
class DocumentSegment:
    """A chunk of text with stable positioning metadata."""
    segment_id: str
    text: str
    start: int = 0  # char offset
    end: int = 0
    page: Optional[int] = None
    timestamp: Optional[float] = None
    segment_type: str = "text"  # text | heading | list_item | table_cell

    def to_dict(self) -> Dict[str, Any]:
        return asdict(self)


@dataclass
class MultimodalDocument:
    """Canonical multimodal artefact returned by the parser.

    All 12 services understand this shape; :mod:`common.multimodal_adapter`
    flattens it onto each service's native request/response contract.
    """
    doc_id: str
    modality: str
    text: str = ""
    segments: List[DocumentSegment] = field(default_factory=list)
    images: List[DocumentImage] = field(default_factory=list)
    tables: List[DocumentTable] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)
    embeddings_ref: Optional[str] = None  # set after /api/v1/multimodal/embed
    warnings: List[str] = field(default_factory=list)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "doc_id": self.doc_id,
            "modality": self.modality,
            "text": self.text,
            "segments": [s.to_dict() for s in self.segments],
            "images": [i.to_dict() for i in self.images],
            "tables": [t.to_dict() for t in self.tables],
            "metadata": dict(self.metadata),
            "embeddings_ref": self.embeddings_ref,
            "warnings": list(self.warnings),
        }

    # Convenience accessors -------------------------------------------------
    @property
    def has_text(self) -> bool:
        return bool(self.text and self.text.strip())

    @property
    def has_images(self) -> bool:
        return bool(self.images)

    @property
    def has_tables(self) -> bool:
        return bool(self.tables)

    def concat_text(self, sep: str = "\n\n") -> str:
        return sep.join(s.text for s in self.segments if s.text)


# ---------------------------------------------------------------------------
# Modality detection
# ---------------------------------------------------------------------------
def detect_modality(source: Union[str, bytes, Dict[str, Any]],
                    filename: Optional[str] = None,
                    mime_type: Optional[str] = None) -> str:
    """Return the canonical modality for the input.

    Args:
        source: file path, raw bytes, or a dict with "modality" key.
        filename: optional filename used to disambiguate by extension.
        mime_type: optional MIME type.
    """
    if isinstance(source, dict) and "modality" in source:
        return source["modality"]
    # Extension wins (most reliable signal in tests and batch imports).
    if filename:
        ext = os.path.splitext(filename)[1].lower()
        if ext in EXT_TO_MODALITY:
            return EXT_TO_MODALITY[ext]
    if isinstance(source, str) and not isinstance(source, bytes):
        # If it looks like a file path, prefer the extension.
        ext = os.path.splitext(source)[1].lower()
        if ext in EXT_TO_MODALITY:
            return EXT_TO_MODALITY[ext]
        # If it's a real file, also try libmagic for the MIME type.
        if os.path.isfile(source):
            try:
                import magic  # type: ignore  # python-magic
                mt = magic.from_file(source, mime=True)
                if mt in MIME_TO_MODALITY:
                    return MIME_TO_MODALITY[mt]
            except Exception:  # noqa: BLE001
                pass
    if mime_type and mime_type in MIME_TO_MODALITY:
        return MIME_TO_MODALITY[mime_type]
    return MODALITY_TEXT


# ---------------------------------------------------------------------------
# Lightweight helpers
# ---------------------------------------------------------------------------
def _new_id(prefix: str = "doc") -> str:
    return f"{prefix}-{uuid.uuid4().hex[:12]}"


def _sha256(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _read_source(source: Union[str, bytes, os.PathLike],
                 filename: Optional[str] = None) -> Tuple[bytes, str, str]:
    """Return (raw_bytes, filename, mime_type)."""
    if isinstance(source, (bytes, bytearray)):
        raw = bytes(source)
        mt = "application/octet-stream"
        if filename:
            ext = os.path.splitext(filename)[1].lower()
            mt = (mimetypes.guess_type(filename)[0] or mt)
        return raw, filename or "blob.bin", mt
    path = str(source)
    if not os.path.isfile(path):
        raise FileNotFoundError(path)
    with open(path, "rb") as fh:
        raw = fh.read()
    fn = filename or os.path.basename(path)
    mt = mimetypes.guess_type(fn)[0] or "application/octet-stream"
    return raw, fn, mt


def _decode_wav_pcm16(raw: bytes) -> Optional[Tuple[int, List[float]]]:
    """Decode WAV PCM16 → (sample_rate, mono_float32_list)."""
    try:
        with wave.open(io.BytesIO(raw), "rb") as w:
            sr = w.getframerate()
            nframes = w.getnframes()
            nchan = w.getnchannels()
            sampw = w.getsampwidth()
            raw_frames = w.readframes(nframes)
        if sampw != 2:
            return None  # only PCM16 for now
        # interleaved channels → mono
        total = len(raw_frames) // 2
        samples = struct.unpack(f"<{total}h", raw_frames)
        if nchan > 1:
            samples = samples[::nchan]
        peak = max(1, max(abs(s) for s in samples))
        return sr, [s / peak for s in samples]
    except Exception as e:  # noqa: BLE001
        logger.debug("wav decode failed: %s", e)
        return None


# ---------------------------------------------------------------------------
# Parser
# ---------------------------------------------------------------------------
class MultiModalParser:
    """Parse any of 6 input modalities into :class:`MultimodalDocument`.

    Optional dependencies degrade gracefully — see ``warnings`` on the
    returned document.  Designed to be safe to instantiate in any thread
    (no shared mutable state besides a single lock).
    """

    SUPPORTED_INPUT_FORMATS: Tuple[str, ...] = (
        "pdf", "docx", "xlsx", "pptx", "md", "html", "txt",
        "jpg", "jpeg", "png", "webp", "gif", "bmp",
        "mp4", "avi", "mov", "mkv", "webm",
        "mp3", "wav", "flac", "m4a", "ogg",
        "eml", "msg", "mbox",
    )

    def __init__(self, *,
                 max_pdf_pages: int = 100,
                 max_excel_rows: int = 5000,
                 max_video_frames: int = 8,
                 max_audio_seconds: float = 60.0,
                 include_image_base64: bool = False) -> None:
        self.max_pdf_pages = max_pdf_pages
        self.max_excel_rows = max_excel_rows
        self.max_video_frames = max_video_frames
        self.max_audio_seconds = max_audio_seconds
        self.include_image_base64 = include_image_base64

    # ----- public entry point ---------------------------------------------
    def parse(self, source: Union[str, bytes, os.PathLike],
              *,
              filename: Optional[str] = None,
              mime_type: Optional[str] = None,
              modality: Optional[str] = None) -> MultimodalDocument:
        raw, fn, mt = _read_source(source, filename=filename)
        if modality is None:
            modality = detect_modality(fn, mime_type=mt)
        doc = MultimodalDocument(
            doc_id=_new_id("mm"),
            modality=modality,
            metadata={
                "filename": fn,
                "mime_type": mt,
                "size_bytes": len(raw),
                "sha256": _sha256(raw),
            },
        )
        try:
            ext = os.path.splitext(fn)[1].lower().lstrip(".")
            if modality == MODALITY_DOCUMENT:
                self._parse_document(doc, raw, ext)
            elif modality == MODALITY_IMAGE:
                self._parse_image(doc, raw, ext)
            elif modality == MODALITY_VIDEO:
                self._parse_video(doc, raw, ext)
            elif modality == MODALITY_AUDIO:
                self._parse_audio(doc, raw, ext)
            elif modality == MODALITY_EMAIL:
                self._parse_email(doc, raw, ext)
            elif modality == MODALITY_TEXT:
                self._parse_text(doc, raw)
            else:  # multimix / unknown
                self._parse_text(doc, raw)
                doc.warnings.append(f"unknown_modality:{modality}")
        except Exception as e:  # noqa: BLE001
            logger.warning("parse failed for %s: %s", fn, e)
            doc.warnings.append(f"parse_error:{e!s}")
            if not doc.text:
                # best effort: at least the raw text decode
                doc.text = raw.decode("utf-8", errors="ignore")[:5000]
        return doc

    def parse_batch(self, sources: Iterable[Union[str, bytes, os.PathLike]],
                    *, filename: Optional[str] = None,
                    mime_type: Optional[str] = None,
                    modality: Optional[str] = None) -> List[MultimodalDocument]:
        return [self.parse(s, filename=filename, mime_type=mime_type,
                          modality=modality) for s in sources]

    # ----- document formats ------------------------------------------------
    def _parse_document(self, doc: MultimodalDocument, raw: bytes, ext: str) -> None:
        if ext == "pdf":
            self._parse_pdf(doc, raw)
        elif ext in ("docx",):
            self._parse_docx(doc, raw)
        elif ext in ("xlsx",):
            self._parse_xlsx(doc, raw)
        elif ext in ("pptx",):
            self._parse_pptx(doc, raw)
        elif ext in ("md", "markdown"):
            self._parse_markdown(doc, raw)
        elif ext in ("html", "htm"):
            self._parse_html(doc, raw)
        elif ext in ("txt",):
            self._parse_text(doc, raw)
        else:
            # generic — try text decode
            self._parse_text(doc, raw)
            doc.warnings.append(f"unknown_doc_ext:{ext}")

    def _parse_pdf(self, doc: MultimodalDocument, raw: bytes) -> None:
        text_chunks: List[str] = []
        n_pages = 0
        # 1) pdfplumber for tables + text
        try:
            import pdfplumber  # type: ignore
            with pdfplumber.open(io.BytesIO(raw)) as pdf:
                doc.metadata["pages"] = len(pdf.pages)
                for i, page in enumerate(pdf.pages[: self.max_pdf_pages]):
                    n_pages += 1
                    t = page.extract_text() or ""
                    if t:
                        text_chunks.append(t)
                        doc.segments.append(DocumentSegment(
                            segment_id=f"seg-p{i+1}",
                            text=t, start=len(doc.text), end=len(doc.text) + len(t),
                            page=i + 1, segment_type="text",
                        ))
                        doc.text += t + "\n\n"
                    # tables
                    try:
                        for ti, tbl in enumerate(page.extract_tables() or []):
                            rows = [[(c or "").strip() for c in row] for row in tbl]
                            doc.tables.append(DocumentTable(
                                table_id=f"tbl-p{i+1}-{ti+1}",
                                rows=rows, page=i + 1,
                            ))
                    except Exception:  # noqa: BLE001
                        pass
                    # images
                    try:
                        for ii, img in enumerate(page.images or []):
                            w = int(img.get("width", 0))
                            h = int(img.get("height", 0))
                            doc.images.append(DocumentImage(
                                image_id=f"img-p{i+1}-{ii+1}",
                                mime_type="image/unknown",
                                width=w, height=h, page=i + 1,
                            ))
                    except Exception:  # noqa: BLE001
                        pass
        except Exception as e:  # noqa: BLE001
            doc.warnings.append(f"pdfplumber_missing:{e!s}")
            # 2) fallback: pypdf for text only
            try:
                from pypdf import PdfReader  # type: ignore
                reader = PdfReader(io.BytesIO(raw))
                doc.metadata["pages"] = len(reader.pages)
                for i, page in enumerate(reader.pages[: self.max_pdf_pages]):
                    n_pages += 1
                    t = page.extract_text() or ""
                    if t:
                        text_chunks.append(t)
                        doc.segments.append(DocumentSegment(
                            segment_id=f"seg-p{i+1}",
                            text=t, start=len(doc.text), end=len(doc.text) + len(t),
                            page=i + 1, segment_type="text",
                        ))
                        doc.text += t + "\n\n"
            except Exception as e2:  # noqa: BLE001
                doc.warnings.append(f"pypdf_missing:{e2!s}")
        doc.metadata["pages_parsed"] = n_pages
        # 3) extract embedded raster images via pymupdf if available
        try:
            import fitz  # type: ignore  # PyMuPDF
            with fitz.open(stream=raw, filetype="pdf") as pdf:
                if "pages" not in doc.metadata:
                    doc.metadata["pages"] = len(pdf)
                for pi, page in enumerate(pdf[: self.max_pdf_pages]):
                    for ii, img_info in enumerate(page.get_images(full=True) or []):
                        xref = img_info[0]
                        try:
                            base = pdf.extract_image(xref)
                        except Exception:  # noqa: BLE001
                            continue
                        b = base.get("image", b"")
                        if not b:
                            continue
                        b64 = base64.b64encode(b).decode("ascii") if self.include_image_base64 else ""
                        doc.images.append(DocumentImage(
                            image_id=f"img-x{pi+1}-{ii+1}",
                            mime_type=base.get("ext", "png") and f"image/{base.get('ext')}",
                            width=int(base.get("width", 0)),
                            height=int(base.get("height", 0)),
                            page=pi + 1,
                            bytes_size=len(b),
                            sha256=_sha256(b),
                            base64=b64,
                        ))
        except Exception:  # noqa: BLE001
            doc.warnings.append("pymupdf_unavailable:embedded_images_skipped")

    def _parse_docx(self, doc: MultimodalDocument, raw: bytes) -> None:
        try:
            from docx import Document  # type: ignore
            d = Document(io.BytesIO(raw))
            for p in d.paragraphs:
                t = (p.text or "").strip()
                if not t:
                    continue
                kind = "heading" if p.style and p.style.name.lower().startswith("heading") else "text"
                doc.segments.append(DocumentSegment(
                    segment_id=_new_id("seg"),
                    text=t, start=len(doc.text), end=len(doc.text) + len(t),
                    segment_type=kind,
                ))
                doc.text += t + "\n\n"
            for ti, tbl in enumerate(d.tables):
                rows = [[cell.text for cell in row.cells] for row in tbl.rows]
                doc.tables.append(DocumentTable(
                    table_id=f"tbl-{ti+1}", rows=rows,
                ))
            # inline images
            for rel in d.part.rels.values():
                if "image" in rel.reltype:
                    blob = rel.target_part.blob
                    doc.images.append(DocumentImage(
                        image_id=_new_id("img"),
                        mime_type=rel.target_part.content_type,
                        bytes_size=len(blob),
                        sha256=_sha256(blob),
                        base64=base64.b64encode(blob).decode("ascii") if self.include_image_base64 else "",
                    ))
            doc.metadata["paragraphs"] = len(d.paragraphs)
            doc.metadata["tables"] = len(d.tables)
            doc.metadata["images"] = len(doc.images)
        except Exception as e:  # noqa: BLE001
            doc.warnings.append(f"docx_parse_failed:{e!s}")
            self._parse_text(doc, raw)

    def _parse_xlsx(self, doc: MultimodalDocument, raw: bytes) -> None:
        try:
            from openpyxl import load_workbook  # type: ignore
            wb = load_workbook(io.BytesIO(raw), data_only=True, read_only=True)
            doc.metadata["sheets"] = wb.sheetnames
            total_rows = 0
            for sname in wb.sheetnames:
                ws = wb[sname]
                rows: List[List[str]] = []
                for ri, row in enumerate(ws.iter_rows(values_only=True)):
                    if total_rows >= self.max_excel_rows:
                        doc.warnings.append("xlsx_truncated:row_cap")
                        break
                    cells = [("" if c is None else str(c)) for c in row]
                    if any(cells):
                        rows.append(cells)
                    total_rows += 1
                if rows:
                    doc.tables.append(DocumentTable(
                        table_id=f"tbl-{sname}",
                        rows=rows, caption=sname,
                    ))
                    doc.text += f"## Sheet: {sname}\n" + "\n".join(
                        "\t".join(r) for r in rows[:50]) + "\n\n"
                    doc.segments.append(DocumentSegment(
                        segment_id=_new_id("seg"),
                        text=f"Sheet {sname} ({len(rows)} rows)",
                        start=len(doc.text), end=len(doc.text),
                        segment_type="table_cell",
                    ))
            wb.close()
        except Exception as e:  # noqa: BLE001
            doc.warnings.append(f"xlsx_parse_failed:{e!s}")
            self._parse_text(doc, raw)

    def _parse_pptx(self, doc: MultimodalDocument, raw: bytes) -> None:
        try:
            from pptx import Presentation  # type: ignore  # type: ignore
            prs = Presentation(io.BytesIO(raw))
            for si, slide in enumerate(prs.slides, 1):
                slide_text: List[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = "".join(run.text for run in para.runs).strip()
                            if t:
                                slide_text.append(t)
                                doc.text += t + "\n"
                    if shape.shape_type == 13:  # picture
                        try:
                            blob = shape.image.blob
                            doc.images.append(DocumentImage(
                                image_id=f"img-s{si}-{len(doc.images)+1}",
                                mime_type=shape.image.content_type,
                                bytes_size=len(blob),
                                sha256=_sha256(blob),
                                base64=base64.b64encode(blob).decode("ascii") if self.include_image_base64 else "",
                            ))
                        except Exception:  # noqa: BLE001
                            pass
                doc.segments.append(DocumentSegment(
                    segment_id=f"seg-s{si}",
                    text="\n".join(slide_text),
                    start=len(doc.text), end=len(doc.text),
                    page=si, segment_type="heading",
                ))
                doc.text += "\n"
            doc.metadata["slides"] = len(prs.slides)
        except Exception as e:  # noqa: BLE001
            doc.warnings.append(f"pptx_parse_failed:{e!s}")
            # try to extract text from the zip directly
            self._extract_text_from_zip(doc, raw, member_pattern=r"slide\d+\.xml$")

    def _parse_markdown(self, doc: MultimodalDocument, raw: bytes) -> None:
        try:
            import markdown as md_lib  # type: ignore
            text = raw.decode("utf-8", errors="ignore")
            doc.text = text
            # split by headers for segments
            for i, chunk in enumerate(re.split(r"\n(?=#+\s)", text)):
                if not chunk.strip():
                    continue
                kind = "heading" if chunk.lstrip().startswith("#") else "text"
                doc.segments.append(DocumentSegment(
                    segment_id=_new_id("seg"),
                    text=chunk, start=0, end=len(chunk),
                    segment_type=kind,
                ))
            try:
                html = md_lib.markdown(text)
                doc.metadata["html_length"] = len(html)
            except Exception:  # noqa: BLE001
                pass
        except Exception as e:  # noqa: BLE001
            doc.warnings.append(f"markdown_failed:{e!s}")
            self._parse_text(doc, raw)

    def _parse_html(self, doc: MultimodalDocument, raw: bytes) -> None:
        try:
            from bs4 import BeautifulSoup  # type: ignore
            soup = BeautifulSoup(raw, "lxml")
            text = soup.get_text("\n", strip=True)
            doc.text = text
            for ti, tbl in enumerate(soup.find_all("table")):
                rows: List[List[str]] = []
                for tr in tbl.find_all("tr"):
                    cells = [td.get_text(" ", strip=True) for td in tr.find_all(["td", "th"])]
                    if cells:
                        rows.append(cells)
                if rows:
                    doc.tables.append(DocumentTable(
                        table_id=f"tbl-{ti+1}", rows=rows,
                    ))
            for ii, img in enumerate(soup.find_all("img")):
                src = img.get("src", "")
                if not src:
                    continue
                doc.images.append(DocumentImage(
                    image_id=f"img-{ii+1}",
                    mime_type="image/unknown",
                    caption=str(img.get("alt", "")),
                ))
            doc.segments.append(DocumentSegment(
                segment_id=_new_id("seg"),
                text=text, start=0, end=len(text), segment_type="text",
            ))
        except Exception as e:  # noqa: BLE001
            doc.warnings.append(f"html_failed:{e!s}")
            self._parse_text(doc, raw)

    def _parse_text(self, doc: MultimodalDocument, raw: bytes) -> None:
        text = raw.decode("utf-8", errors="ignore")
        doc.text = text
        # simple paragraph split
        for chunk in re.split(r"\n\s*\n", text):
            if not chunk.strip():
                continue
            doc.segments.append(DocumentSegment(
                segment_id=_new_id("seg"),
                text=chunk, start=0, end=len(chunk), segment_type="text",
            ))

    def _extract_text_from_zip(self, doc: MultimodalDocument, raw: bytes,
                               member_pattern: str) -> None:
        try:
            pat = re.compile(member_pattern)
            with zipfile.ZipFile(io.BytesIO(raw)) as z:
                for name in z.namelist():
                    if pat.search(name):
                        try:
                            t = z.read(name).decode("utf-8", errors="ignore")
                            # crude tag strip
                            t = re.sub(r"<[^>]+>", " ", t)
                            t = re.sub(r"\s+", " ", t).strip()
                            if t:
                                doc.text += t + "\n"
                        except Exception:  # noqa: BLE001
                            continue
        except Exception as e:  # noqa: BLE001
            doc.warnings.append(f"zip_extract_failed:{e!s}")

    # ----- images ---------------------------------------------------------
    def _parse_image(self, doc: MultimodalDocument, raw: bytes, ext: str) -> None:
        try:
            from PIL import Image  # type: ignore
            img = Image.open(io.BytesIO(raw))
            img.verify()  # raises if corrupt
            img2 = Image.open(io.BytesIO(raw))  # need a fresh handle
            w, h = img2.size
            doc.images.append(DocumentImage(
                image_id=_new_id("img"),
                mime_type=f"image/{ext if ext != 'jpg' else 'jpeg'}",
                width=w, height=h,
                bytes_size=len(raw), sha256=_sha256(raw),
                base64=base64.b64encode(raw).decode("ascii") if self.include_image_base64 else "",
            ))
            doc.metadata["dimensions"] = [w, h]
            doc.metadata["format"] = img2.format
        except Exception as e:  # noqa: BLE001
            doc.warnings.append(f"image_verify_failed:{e!s}")
        # OCR via pytesseract when available
        try:
            import pytesseract  # type: ignore
            from PIL import Image  # type: ignore
            txt = pytesseract.image_to_string(Image.open(io.BytesIO(raw)))
            txt = (txt or "").strip()
            if txt:
                doc.text = txt
                doc.segments.append(DocumentSegment(
                    segment_id=_new_id("ocr"),
                    text=txt, start=0, end=len(txt),
                    segment_type="text",
                ))
                if doc.images:
                    doc.images[0].ocr_text = txt
        except Exception:  # noqa: BLE001
            doc.warnings.append("ocr_unavailable:pytesseract_missing")

    # ----- video ---------------------------------------------------------
    def _parse_video(self, doc: MultimodalDocument, raw: bytes, ext: str) -> None:
        doc.metadata["container"] = ext
        # Probe via OpenCV if available — extract keyframes as images
        try:
            import cv2  # type: ignore
            import numpy as np  # type: ignore
            with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
                tmp.write(raw)
                tmp_path = tmp.name
            try:
                cap = cv2.VideoCapture(tmp_path)
                if cap.isOpened():
                    fps = cap.get(cv2.CAP_PROP_FPS) or 0.0
                    total = int(cap.get(cv2.CAP_PROP_FRAME_COUNT) or 0)
                    width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH) or 0)
                    height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT) or 0)
                    doc.metadata["fps"] = float(fps)
                    doc.metadata["frames"] = total
                    doc.metadata["width"] = width
                    doc.metadata["height"] = height
                    doc.metadata["duration_s"] = float(total / fps) if fps else 0.0
                    if total > 0:
                        step = max(1, total // self.max_video_frames)
                        idx = 0
                        kept = 0
                        while kept < self.max_video_frames and idx < total:
                            cap.set(cv2.CAP_PROP_POS_FRAMES, idx)
                            ok, frame = cap.read()
                            if not ok:
                                break
                            ts = idx / fps if fps else 0.0
                            ok2, buf = cv2.imencode(".jpg", frame)
                            if not ok2:
                                idx += step
                                continue
                            blob = buf.tobytes()
                            doc.images.append(DocumentImage(
                                image_id=f"frame-{idx}",
                                mime_type="image/jpeg",
                                width=int(frame.shape[1]),
                                height=int(frame.shape[0]),
                                timestamp=round(ts, 3),
                                bytes_size=len(blob),
                                sha256=_sha256(blob),
                                base64=base64.b64encode(blob).decode("ascii") if self.include_image_base64 else "",
                            ))
                            doc.segments.append(DocumentSegment(
                                segment_id=f"seg-t{idx}",
                                text=f"[video frame @ {ts:.2f}s]",
                                start=idx, end=idx,
                                timestamp=round(ts, 3),
                                segment_type="text",
                            ))
                            idx += step
                            kept += 1
                cap.release()
            finally:
                try:
                    os.unlink(tmp_path)
                except Exception:  # noqa: BLE001
                    pass
        except Exception as e:  # noqa: BLE001
            doc.warnings.append(f"video_probe_failed:{e!s}")
            doc.text = f"[video:{ext} bytes={len(raw)}]"

    # ----- audio ---------------------------------------------------------
    def _parse_audio(self, doc: MultimodalDocument, raw: bytes, ext: str) -> None:
        doc.metadata["container"] = ext
        # 1) Try WAV PCM16 decode for fallback features
        wav = _decode_wav_pcm16(raw) if ext == "wav" else None
        if wav is not None:
            sr, samples = wav
            doc.metadata["sample_rate"] = sr
            doc.metadata["samples"] = len(samples)
            doc.metadata["duration_s"] = round(len(samples) / sr, 3)
        else:
            # approximate duration by bitrate
            est_dur = len(raw) / 16000.0  # 128 kbps mp3 ~ 16KB/s
            doc.metadata["duration_s_estimated"] = round(est_dur, 3)
        # 2) ASR via whisper
        try:
            import whisper  # type: ignore
            with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
                tmp.write(raw)
                tmp_path = tmp.name
            try:
                # prefer the tiny model for offline robustness
                try:
                    model = whisper.load_model("tiny")
                except Exception:  # noqa: BLE001
                    model = None
                if model is not None:
                    res = model.transcribe(tmp_path, fp16=False)
                    text = (res.get("text") or "").strip()
                    if text:
                        doc.text = text
                        for seg in res.get("segments") or []:
                            doc.segments.append(DocumentSegment(
                                segment_id=_new_id("asr"),
                                text=seg.get("text", "").strip(),
                                start=int(seg.get("start", 0) * 1000),
                                end=int(seg.get("end", 0) * 1000),
                                timestamp=round(seg.get("start", 0.0), 3),
                                segment_type="text",
                            ))
            finally:
                try:
                    os.unlink(tmp_path)
                except Exception:  # noqa: BLE001
                    pass
        except Exception as e:  # noqa: BLE001
            doc.warnings.append(f"whisper_unavailable:{e!s}")

    # ----- email ---------------------------------------------------------
    def _parse_email(self, doc: MultimodalDocument, raw: bytes, ext: str) -> None:
        try:
            msg = email.message_from_bytes(raw, policy=email.policy.default)
            doc.metadata["subject"] = str(msg.get("Subject", ""))
            doc.metadata["from"] = str(msg.get("From", ""))
            doc.metadata["to"] = str(msg.get("To", ""))
            doc.metadata["date"] = str(msg.get("Date", ""))
            body_parts: List[str] = []
            for part in msg.walk():
                ctype = part.get_content_type()
                if ctype == "text/plain" and not part.is_multipart():
                    body_parts.append(part.get_content())
                elif ctype == "text/html" and not part.is_multipart():
                    try:
                        from bs4 import BeautifulSoup  # type: ignore
                        body_parts.append(
                            BeautifulSoup(part.get_content(), "lxml").get_text("\n", strip=True)
                        )
                    except Exception:  # noqa: BLE001
                        body_parts.append(part.get_content())
                elif part.get_filename():
                    blob = part.get_payload(decode=True) or b""
                    doc.images.append(DocumentImage(
                        image_id=_new_id("att"),
                        mime_type=ctype,
                        bytes_size=len(blob),
                        sha256=_sha256(blob),
                        caption=part.get_filename(),
                        base64=base64.b64encode(blob).decode("ascii") if self.include_image_base64 else "",
                    ))
            doc.text = "\n\n".join(p for p in body_parts if p)
            doc.segments.append(DocumentSegment(
                segment_id=_new_id("body"),
                text=doc.text, start=0, end=len(doc.text), segment_type="text",
            ))
        except Exception as e:  # noqa: BLE001
            doc.warnings.append(f"email_parse_failed:{e!s}")
            self._parse_text(doc, raw)


__all__ = [
    "MultiModalParser",
    "MultimodalDocument",
    "DocumentImage",
    "DocumentTable",
    "DocumentSegment",
    "detect_modality",
    "MODALITY_TEXT", "MODALITY_IMAGE", "MODALITY_VIDEO",
    "MODALITY_AUDIO", "MODALITY_DOCUMENT", "MODALITY_EMAIL",
    "MODALITY_MULTIMIX",
    "ALL_MODALITIES",
    "OUTPUT_TEXT", "OUTPUT_JSON", "OUTPUT_MULTIMODAL", "ALL_OUTPUT_KINDS",
    "MIME_TO_MODALITY", "EXT_TO_MODALITY",
]
